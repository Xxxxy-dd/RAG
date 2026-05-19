"""PPTX loader to Markdown format.

PPT 转 Markdown 的核心逻辑：
- 标题提取：优先识别每一页的主标题，转化为 Markdown 的标题（## 标题）。
- 文本框识别：同一页里的多文本框按阅读顺序（X/Y坐标）排序。
- 列表识别：保留项目符号，转换为标准的 Markdown 无序列表（- 列表项）。
- 表格提取：识别 PPT 表格对象，转化为 Markdown 表格语法（| 列1 | 列2 |）。
- 噪声过滤：去除页码、模板导航词（如 PART 01）。

截断工作应交由后续通用的 MarkdownTextSplitter 完成。
"""

import re
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_core.documents import Document


_SLIDE_NUMBER_RE = re.compile(r"^(第\s*\d+\s*页|\d+\s*/\s*\d+|page\s*\d+|\d+)$", re.IGNORECASE)
_BULLET_RE = re.compile(r"^([\-\*•·]|\d+[\.)]|[一二三四五六七八九十]+[、\.)]|\([一二三四五六七八九十]+\))\s*")
_NOISE_WORDS = {"内容", "program", "part", "目录", "agenda", "contents", "thank you", "q&a"}


def _normalize_text(text: str) -> str:
    """把空格、换行和多余的空白整理干净。"""
    if not text:
        return ""

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _is_noise_text(text: str) -> bool:
    """判断一段文字是不是页码、装饰词、大纲导航或者模板噪声。"""
    if not text:
        return True

    cleaned = _normalize_text(text).strip()
    if not cleaned:
        return True

    if _SLIDE_NUMBER_RE.match(cleaned):
        return True

    if len(cleaned) <= 3 and not re.search(r"[\u4e00-\u9fffA-Za-z0-9]", cleaned):
        return True

    # 判断是否为固定的模板导航词
    lower_cleaned = cleaned.lower()
    core_text = re.sub(r'[\d\.\-\*•· ]+', '', lower_cleaned)
    if core_text in _NOISE_WORDS or any(word == lower_cleaned for word in _NOISE_WORDS):
        return True
        
    if re.match(r"^part\s*\d+(\s*[•·\-\*])?$", lower_cleaned):
        return True

    return False


def _shape_position(shape):
    """拿到元素在页面上的位置，用于按阅读顺序排版。"""
    left = getattr(shape, "left", 0) or 0
    top = getattr(shape, "top", 0) or 0
    width = getattr(shape, "width", 0) or 0
    height = getattr(shape, "height", 0) or 0
    return top, left, width, height


def _shape_reading_key(shape):
    """给每个元素算一个阅读顺序，标题优先，再按从上到下、从左到右。"""
    top, left, width, height = _shape_position(shape)
    is_title = 0
    if getattr(shape, "is_placeholder", False):
        placeholder = shape.placeholder_format
        if placeholder and placeholder.type == 1:  # TITLE
            is_title = -10

    return is_title, top, left, -width, -height


def _paragraph_text(paragraph) -> str:
    """提取段落文字，并将 PPT 面板里的项目符号统一转成 Markdown 列表格式 `- `"""
    text = _normalize_text(paragraph.text)
    if not text:
        return ""

    if _BULLET_RE.match(text):
        # 替换开头的各种奇异项目符号为标准的短横线
        text = _BULLET_RE.sub("- ", text)
        return text

    if getattr(paragraph, "level", 0) and paragraph.level > 0:
        indent = "  " * paragraph.level
        return f"{indent}- {text}"

    return text


def _get_shape_text(shape) -> str:
    """提取幻灯片元素的内容。如果遇到表格，将其直接转化为 Markdown 表格。"""
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        rows_md = []
        table = shape.table
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                # 剔除单元格内部换行，防止破坏 Markdown 表格结构
                cell_text = _normalize_text(cell.text).replace("\n", " ")
                cells.append(cell_text)
                
            # 拼接 Markdown 行
            rows_md.append("| " + " | ".join(cells) + " |")
            
            # 第一行结束后，添加 Markdown 的表头分割线
            if i == 0:
                separator = "|" + "|".join(["---"] * len(cells)) + "|"
                rows_md.append(separator)
                
        return "\n".join(rows_md)

    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        lines = []
        for paragraph in shape.text_frame.paragraphs:
            line = _paragraph_text(paragraph)
            if line:
                lines.append(line)
        return "\n".join(lines)

    return ""


def _extract_slide_title(slide) -> str:
    """提取页面的大标题"""
    # 优先找标准的标题占位符
    for shape in sorted(slide.shapes, key=_shape_reading_key):
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue

        if getattr(shape, "is_placeholder", False):
            placeholder = shape.placeholder_format
            if placeholder and placeholder.type == 1:
                title = _normalize_text(shape.text)
                if title:
                    return title

    # 如果没找到，就拿顶部的、没有很长的第一行文字兜底
    for shape in sorted(slide.shapes, key=_shape_reading_key):
        text = _get_shape_text(shape)
        if text and len(text) <= 60 and not _is_noise_text(text):
            first_line = text.split("\n")[0].strip()
            if first_line:
                return first_line

    return ""


def load_pptx(path: str) -> List[Document]:
    """主入口：将 PPT 文件完整转化为包含结构化信息的 Markdown Document 列表。"""
    prs = Presentation(path)
    docs = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        title = _extract_slide_title(slide)
        blocks = []
        
        # 1. 加入 Markdown 级别的标题
        if title:
            blocks.append(f"## {title}")
            
        # 2. 按顺序提取内部所有文本框和表格
        for shape in sorted(slide.shapes, key=_shape_reading_key):
            text = _get_shape_text(shape)
            if not text:
                continue

            # 去重：跳过前面已经提取为大标题的内容
            if title and text.strip() == title:
                continue

            # 过滤噪声
            filtered_lines = []
            for line in text.split("\n"):
                if not _is_noise_text(line) and line.strip() != title:
                    filtered_lines.append(line)
                    
            cleaned_text = "\n".join(filtered_lines).strip()
            if cleaned_text:
                blocks.append(cleaned_text)

        # 3. 拼接整页的 Markdown
        md_content = "\n\n".join(blocks).strip()
        if not md_content:
            continue

        # 4. 保留元数据
        metadata = {
            "source": path,
            "slide_index": slide_index,
            "slide_title": title or f"Slide {slide_index}",
        }
        
        docs.append(Document(page_content=md_content, metadata=metadata))

    # 后处理：合并过短的导航页。
    # 比如仅包含“第一部分”这类文字的页面，内容极短，单独成块没有意义。
    merged_docs = []
    for doc in docs:
        content_length = len(doc.page_content.strip())
        # 如果这是过渡页（通常字数非常少，例如只包含“## 第一部分”或类似于“第X章 介绍”），并且上一个/下一个文档可以合并
        if content_length < 50:
            if not merged_docs:
                merged_docs.append(doc)
            else:
                # 策略：如果内容太短，就往它前面的文档合并（作为标题或过渡），并将元数据的标题合并
                prev_doc = merged_docs[-1]
                prev_doc.page_content += "\n\n" + doc.page_content
                # 稍微保留一下跨页信息
                prev_doc.metadata["merged_slides"] = prev_doc.metadata.get("merged_slides", [prev_doc.metadata["slide_index"]]) + [doc.metadata["slide_index"]]
        else:
            # 或者，如果上一个非常短，把它补充到当前的开头（这取决于导航页通常是引导下一页）
            if merged_docs and len(merged_docs[-1].page_content.strip()) < 50:
                short_doc = merged_docs.pop()
                doc.page_content = short_doc.page_content + "\n\n" + doc.page_content
                doc.metadata["merged_slides"] = [short_doc.metadata["slide_index"], doc.metadata["slide_index"]]
            merged_docs.append(doc)

    return merged_docs



